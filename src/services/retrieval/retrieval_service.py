"""Retrieval helpers for evidence search and template deck selection."""

from pathlib import Path

from pptx import Presentation

from core.config import Settings
from core.logger import get_logger
from services.retrieval.vector_store import build_or_load_vector_store

logger = get_logger(__name__)


def retrieve_context(query: str, settings: Settings) -> tuple[list[dict], list[str]]:
    """Retrieve the most relevant supporting context for the current RFP."""
    logger.info("Retrieving context for query (top_k=%d)", settings.top_k_results)
    vector_store = build_or_load_vector_store(settings)
    docs = vector_store.similarity_search(query, k=settings.top_k_results)
    context = []
    citations = []

    for doc in docs:
        source = doc.metadata.get("source", "unknown")
        slide_number = doc.metadata.get("slide_number")
        slide_title = doc.metadata.get("slide_title")
        citation = f"{source} (chunk {doc.metadata.get('chunk_id', 'n/a')})"
        if slide_number:
            citation = f"{source} (slide {slide_number}: {slide_title})"
        citations.append(citation)
        # Keep raw text plus metadata so later steps can ground prompts and trace output.
        context.append(
            {
                "text": doc.page_content,
                "source": source,
                "metadata": doc.metadata,
            }
        )

    logger.info("Retrieved %d context chunks, %d citations", len(context), len(citations))
    return context, citations


def retrieve_template_outline(query: str, settings: Settings) -> tuple[str, list[dict]]:
    """Select the best reference deck and rebuild its ordered slide outline."""
    preferred_template = settings.preferred_template_path
    if preferred_template.exists():
        logger.info("Using preferred template: %s", preferred_template)
        template_slides = _load_template_outline_from_ppt(preferred_template)
        if template_slides:
            return str(preferred_template), template_slides

    vector_store = build_or_load_vector_store(settings)
    docs_with_scores = vector_store.similarity_search_with_score(
        query,
        k=max(settings.top_k_results * 3, 12),
    )

    source_rank: dict[str, dict] = {}
    for doc, score in docs_with_scores:
        metadata = doc.metadata
        if metadata.get("content_type") != "slide":
            continue
        source = metadata.get("source")
        if not source:
            continue

        payload = source_rank.setdefault(source, {"hits": 0, "best_score": score})
        payload["hits"] += 1
        payload["best_score"] = min(payload["best_score"], score)

    if not source_rank:
        logger.warning("No slide-type documents found in vector store for template selection")
        return "", []

    selected_source = sorted(
        source_rank.items(),
        key=lambda item: (-item[1]["hits"], item[1]["best_score"]),
    )[0][0]
    logger.info("Selected template source: %s", selected_source)

    template_slides = []
    docstore = getattr(vector_store, "docstore", None)
    records = getattr(docstore, "_dict", {}) if docstore else {}

    for record in records.values():
        metadata = getattr(record, "metadata", {})
        if metadata.get("source") != selected_source:
            continue
        if metadata.get("content_type") != "slide":
            continue

        # Reconstruct the ordered template outline from all slides in the winning deck.
        template_slides.append(
            {
                "title": (
                    metadata.get("slide_title")
                    or f"Slide {metadata.get('slide_number', '')}".strip()
                ),
                "slide_number": metadata.get("slide_number", 0),
                "reference_text": record.page_content,
                "source": selected_source,
            }
        )

    template_slides.sort(key=lambda item: item["slide_number"])
    return selected_source, template_slides


def _load_template_outline_from_ppt(template_path: Path) -> list[dict]:
    """Read slide order, titles, and text directly from a preferred PPT template."""
    presentation = Presentation(str(template_path))
    template_slides: list[dict] = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_title = _extract_slide_title(slide, slide_number)
        reference_text = _extract_slide_text(slide)
        template_slides.append(
            {
                "title": slide_title,
                "slide_number": slide_number,
                "reference_text": reference_text,
                "source": str(template_path),
            }
        )

    return template_slides


def _extract_slide_title(slide, slide_number: int) -> str:
    """Return a stable title for a slide, falling back to the first text block."""
    if slide.shapes.title and slide.shapes.title.text.strip():
        return slide.shapes.title.text.strip()

    for shape in slide.shapes:
        if hasattr(shape, "text_frame"):
            text = " ".join(
                paragraph.text.strip()
                for paragraph in shape.text_frame.paragraphs
                if paragraph.text.strip()
            )
            if text:
                return text[:80]

    return f"Slide {slide_number}"


def _extract_slide_text(slide) -> str:
    """Collect visible text from a slide for reference-aware generation."""
    parts: list[str] = []
    for shape in slide.shapes:
        if not hasattr(shape, "text_frame"):
            continue
        for paragraph in shape.text_frame.paragraphs:
            text = paragraph.text.strip()
            if text:
                parts.append(text)
    return "\n".join(parts)
