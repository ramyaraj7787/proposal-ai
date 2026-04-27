"""Document loading utilities for proposal assets."""

from pathlib import Path

from docx import Document as WordDocument
from pypdf import PdfReader
from pptx import Presentation

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".txt"}


def load_text_from_file(file_path: str | Path) -> str:
    """Load plain text from a supported proposal asset."""
    path = Path(file_path)
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return _load_pdf(path)
    if suffix == ".docx":
        return _load_docx(path)
    if suffix == ".pptx":
        return _load_pptx(path)
    if suffix == ".txt":
        return path.read_text(encoding="utf-8")
    raise ValueError(f"Unsupported file type: {path.suffix}")


def load_documents_from_directory(directory: str | Path) -> list[dict]:
    """Read and normalize all supported documents from a directory tree."""
    root = Path(directory)
    documents: list[dict] = []

    for path in root.rglob("*"):
        if not path.is_file() or path.suffix.lower() not in SUPPORTED_EXTENSIONS:
            continue

        if path.suffix.lower() == ".pptx":
            # PowerPoint files are indexed per slide so we can preserve template order.
            documents.extend(_load_pptx_documents(path))
            continue

        text = load_text_from_file(path)
        if text.strip():
            documents.append(
                {
                    "text": text,
                    "source": str(path),
                    "metadata": {
                        "filename": path.name,
                        "extension": path.suffix.lower(),
                        "parent_folder": path.parent.name,
                        "content_type": "document",
                    },
                }
            )

    return documents


def _load_pdf(path: Path) -> str:
    """Extract concatenated text from a PDF file."""
    reader = PdfReader(str(path))
    return "\n".join(page.extract_text() or "" for page in reader.pages)


def _load_docx(path: Path) -> str:
    """Extract non-empty paragraphs from a Word document."""
    doc = WordDocument(str(path))
    paragraphs = [paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()]
    return "\n".join(paragraphs)


def _load_pptx(path: Path) -> str:
    """Return a whole-deck text representation of a PowerPoint file."""
    return "\n\n".join(item["text"] for item in _load_pptx_documents(path))


def _load_pptx_documents(path: Path) -> list[dict]:
    """Split a PowerPoint deck into slide-level retrieval documents."""
    presentation = Presentation(str(path))
    documents: list[dict] = []

    for slide_idx, slide in enumerate(presentation.slides, start=1):
        slide_lines: list[str] = []
        title = ""

        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if not text:
                continue

            normalized = text.strip()
            if not normalized:
                continue

            if not title:
                # The first non-empty text block acts as the slide title hint.
                title = normalized.splitlines()[0]

            slide_lines.append(normalized)

        if slide_lines:
            # Each slide becomes its own retrieval record with template metadata.
            documents.append(
                {
                    "text": "\n".join(slide_lines),
                    "source": str(path),
                    "metadata": {
                        "filename": path.name,
                        "extension": path.suffix.lower(),
                        "parent_folder": path.parent.name,
                        "content_type": "slide",
                        "slide_number": slide_idx,
                        "slide_title": title or f"Slide {slide_idx}",
                    },
                }
            )

    return documents
