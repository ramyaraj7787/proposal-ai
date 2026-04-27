"""PowerPoint rendering helpers for generated proposal content."""

from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from core.config import Settings

MAX_BULLETS_PER_SLIDE = 3
MAX_BULLET_CHARS = 150
BASE_FONT_SIZE = 18
MIN_FONT_SIZE = 14


def build_proposal_ppt(state: dict, settings: Settings) -> str:
    """Create a proposal deck using the reference PPT as the base template."""
    presentation = _load_base_presentation(state)
    _populate_cover_slide(presentation, state)

    content_template_slides = _get_content_template_slides(state)
    for section, template_slide in zip(state["proposal_sections"], content_template_slides):
        target_slide = _resolve_target_slide(
            presentation,
            template_slide.get("slide_number", 0) - 1,
        )
        bullets = [
            line.strip()
            for line in section["content"].splitlines()
            if line.strip()
        ]
        _populate_content_slide(target_slide, section["title"], bullets)

    output_file = (
        settings.output_dir
        / f"proposal_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    )
    presentation.save(str(output_file))
    return str(output_file)


def _load_base_presentation(state: dict) -> Presentation:
    """Load the selected template deck or fall back to a blank presentation."""
    template_source = state.get("template_source")
    if template_source and Path(template_source).exists():
        return Presentation(template_source)
    return Presentation()


def _populate_cover_slide(presentation: Presentation, state: dict) -> None:
    """Treat the first slide of the reference deck as the proposal cover slide."""
    if not presentation.slides:
        slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    else:
        slide = presentation.slides[0]

    title_shape = slide.shapes.title or _find_first_text_shape(slide)
    if title_shape is not None:
        title_shape.text = "Proposal"

    subtitle_shape = _find_cover_subtitle_shape(slide, title_shape)
    subtitle_text = (
        f"{state['client']}\n"
        f"{state['sector']} | {state['domain']} | {state['country']}\n"
        f"{datetime.now().strftime('%d-%b-%Y')}"
    )
    if subtitle_shape is not None:
        subtitle_shape.text = subtitle_text


def _get_content_template_slides(state: dict) -> list[dict]:
    """Return all content template slides except the cover slide."""
    template_slides = state.get("template_slides", [])
    if not template_slides:
        return []

    first_slide_number = min(
        slide.get("slide_number", 1)
        for slide in template_slides
    )
    return [
        slide
        for slide in template_slides
        if slide.get("slide_number", 1) != first_slide_number
    ]


def _populate_content_slide(slide, title: str, bullets: list[str]) -> None:
    """Write generated content into an existing template slide."""
    if slide.shapes.title is not None:
        slide.shapes.title.text = title

    normalized_bullets = _normalize_bullets(bullets)[:MAX_BULLETS_PER_SLIDE]
    if not normalized_bullets:
        normalized_bullets = ["Content generated from the RFP and retrieved proposal assets."]

    body_shape = _find_body_shape(slide)
    if body_shape is None:
        body_shape = slide.shapes.add_textbox(
            Inches(1.0),
            Inches(1.8),
            Inches(8.0),
            Inches(4.0),
        )

    _write_bullets(body_shape.text_frame, normalized_bullets)


def _resolve_target_slide(presentation: Presentation, slide_index: int):
    """Return the target slide by sequential position or create a fallback slide."""
    if 0 <= slide_index < len(presentation.slides):
        return presentation.slides[slide_index]

    # If the template has fewer slides than expected, append a fallback.
    if len(presentation.slide_layouts) > 1:
        return presentation.slides.add_slide(presentation.slide_layouts[1])
    return presentation.slides.add_slide(presentation.slide_layouts[0])


def _find_cover_subtitle_shape(slide, title_shape):
    """Return the most likely subtitle placeholder on the cover slide."""
    for shape in slide.shapes:
        if shape == title_shape:
            continue
        if getattr(shape, "is_placeholder", False):
            placeholder_type = shape.placeholder_format.type
            if placeholder_type == PP_PLACEHOLDER.SUBTITLE:
                return shape

    for shape in slide.shapes:
        if shape != title_shape and hasattr(shape, "text_frame"):
            return shape
    return None


def _find_body_shape(slide):
    """Find the most appropriate body placeholder or text box for slide content."""
    preferred_types = {
        PP_PLACEHOLDER.BODY,
        PP_PLACEHOLDER.OBJECT,
    }

    for shape in slide.shapes:
        if not getattr(shape, "is_placeholder", False):
            continue
        if shape == slide.shapes.title:
            continue
        if (
            hasattr(shape, "text_frame")
            and shape.placeholder_format.type in preferred_types
        ):
            return shape

    candidate = None
    candidate_area = 0
    for shape in slide.shapes:
        if shape == slide.shapes.title or not hasattr(shape, "text_frame"):
            continue
        area = getattr(shape, "width", 0) * getattr(shape, "height", 0)
        if area > candidate_area:
            candidate_area = area
            candidate = shape
    return candidate


def _find_first_text_shape(slide):
    """Return the first text-capable shape on a slide."""
    for shape in slide.shapes:
        if hasattr(shape, "text_frame"):
            return shape
    return None


def _write_bullets(text_frame, bullets: list[str]) -> None:
    """Write formatted bullets into a text frame while preserving slide layout."""
    text_frame.clear()
    text_frame.word_wrap = True
    font_size = _font_size_for_bullets(bullets)

    for idx, bullet in enumerate(bullets):
        paragraph = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.space_after = Pt(10)
        for run in paragraph.runs:
            run.font.size = Pt(font_size)


def _normalize_bullets(bullets: list[str]) -> list[str]:
    """Clean bullets and keep them short enough for the template slides."""
    normalized: list[str] = []
    for bullet in bullets:
        cleaned = bullet.replace("\u2022", " ")
        cleaned = cleaned.replace("â€œ", '"').replace("â€", '"')
        cleaned = cleaned.replace("â€™", "'").replace("&amp;", "&")
        cleaned = " ".join(cleaned.split())
        if not cleaned:
            continue
        if cleaned[-1] not in ".!?":
            cleaned += "."
        if len(cleaned) > MAX_BULLET_CHARS:
            cleaned = cleaned[: MAX_BULLET_CHARS - 3].rstrip() + "..."
        normalized.append(cleaned)
    return normalized


def _font_size_for_bullets(bullets: list[str]) -> int:
    """Choose a readable font size based on bullet density."""
    longest = max((len(bullet) for bullet in bullets), default=0)
    if longest > 130:
        return MIN_FONT_SIZE
    if longest > 95:
        return 16
    return BASE_FONT_SIZE
